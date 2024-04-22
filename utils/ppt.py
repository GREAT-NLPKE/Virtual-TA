import pptx
from pptx import Presentation
from pptx.enum.shapes import  MSO_SHAPE_TYPE
from PIL import Image
import cv2 as cv
import numpy as np
from copy import deepcopy
from .iterable import Iterable
from .latex import LatexImage
from .ocr_ import ppstructure 

class PPT(Iterable):
    def __init__(self,path,max_length,*args,**kw) -> None:
        super().__init__()
        self.path = path
        self.max_length = max_length
        self.prs = pptx.Presentation(path)
        self.latex = LatexImage('/nvme01/lmj/virtual-ta/latexocr')
        
    def __len__(self):
        return len(self.prs.slides) 
        
    def __getitem__(self,index):
        contexts=[]
        pictures=[]
        slide = self.prs.slides[index]
        res=''
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    if paragraph.text and paragraph.text[-1] not in ['。', '，','.', '！', '？', '；', '…',',']:
                        res = res+paragraph.text+';'
                    else:
                        res = res+paragraph.text   
            if shape.shape_type==MSO_SHAPE_TYPE.PICTURE:
                image = shape.image.blob
                image = cv.imdecode(np.frombuffer(image,np.uint8),cv.IMREAD_COLOR)
                rgb_img = cv.cvtColor(image,cv.COLOR_BGR2RGB)
                image = Image.fromarray(rgb_img)
                structs = ppstructure(image)
                for struct in structs:
                    if struct['type']=='equation':
                        latex = self.latex.img2latex(image.crop(struct['bbox']))
                        if not latex.startswith(r'\begin'):
                            res = res+latex+';'
                    elif struct['type']=='figure':
                        pictures.append(image)
        
        contexts.append(res)
        return self.clip_max_length(contexts), pictures
    
    def clip_max_length(self,sentences:list[str]):
        result=[]
        for sentence in sentences:
            if len(sentence)>self.max_length:
                r=''
                s = deepcopy(sentence)
                while s:
                    i=[s.index(signal) for signal in ',.;，。；' if signal in s]
                    i.append(len(s)-1)
                    i = min(i)
                    if len(r)+i+1 <= self.max_length:
                        r += s[:i+1]
                        s = s[i+1:] if i+1 != len(s) else []
                    else:
                        result.append(r)
                        r=''
                result.append(r)
            else:
                result.append(sentence)
        return result
    
if __name__=='__main__':
    ppt = PPT(path='/nvme01/lmj/virtual-ta/data/数图习题2023.pptx',max_length=200)
    # print(len(ppt))
    for c,p in ppt:
        print(c,p)
    # print(ppt[3])